#!/usr/bin/env python3
"""Builds the KLA PS01 idea-submission deck (9 slides, 16:9, dark theme)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- theme
NAVY = RGBColor(0x0A, 0x19, 0x30)        # slide background
PANEL = RGBColor(0x13, 0x2C, 0x4E)       # card background
PANEL_2 = RGBColor(0x0F, 0x24, 0x40)     # darker card
CYAN = RGBColor(0x37, 0xC8, 0xF0)        # primary accent
LIME = RGBColor(0xC6, 0xE8, 0x4B)        # secondary accent (SEMICON mesh green)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xB7, 0xC6, 0xDA)       # secondary text
FAINT = RGBColor(0x6E, 0x84, 0x9E)       # tertiary text
DANGER = RGBColor(0xF2, 0x6D, 0x6D)

TITLE_FONT = "Avenir Next"
BODY_FONT = "Helvetica Neue"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

# ---------------------------------------------------------------- helpers

def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def _no_line(shape):
    shape.line.fill.background()
    shape.shadow.inherit = False


def rect(slide, x, y, w, h, color, rounded=False, line_color=None, line_w=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if rounded:
        try:
            shp.adjustments[0] = 0.08
        except Exception:
            pass
    if line_color is not None:
        shp.line.color.rgb = line_color
        shp.line.width = line_w or Pt(1)
        shp.shadow.inherit = False
    else:
        _no_line(shp)
    return shp


def text(slide, x, y, w, h, runs, size=14, color=WHITE, bold=False, font=BODY_FONT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0, space_after=4):
    """runs: str | list[str] | list[(text, dict)] per paragraph."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [runs]
    for i, item in enumerate(runs):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = line_spacing
        para.space_after = Pt(space_after)
        if isinstance(item, tuple):
            content, opts = item
        else:
            content, opts = item, {}
        r = para.add_run()
        r.text = content
        f = r.font
        f.size = Pt(opts.get("size", size))
        f.bold = opts.get("bold", bold)
        f.name = opts.get("font", font)
        f.color.rgb = opts.get("color", color)
        if opts.get("italic"):
            f.italic = True
    return tb


def header(slide, kicker, title, page_no):
    rect(slide, Inches(0.55), Inches(0.42), Inches(0.6), Inches(0.055), CYAN)
    text(slide, Inches(0.55), Inches(0.52), Inches(9), Inches(0.3),
         kicker.upper(), size=11, color=CYAN, bold=True, font=BODY_FONT)
    text(slide, Inches(0.55), Inches(0.78), Inches(11.2), Inches(0.75),
         title, size=29, color=WHITE, bold=True, font=TITLE_FONT)
    text(slide, Inches(11.9), Inches(0.5), Inches(0.9), Inches(0.3),
         f"{page_no:02d} / 09", size=10, color=FAINT, align=PP_ALIGN.RIGHT)
    text(slide, Inches(0.55), Inches(7.06), Inches(12.2), Inches(0.3),
         "SEMICON India Hackathon 2026   ·   KLA Track 1 — AI-Based Restoration of Degraded Images",
         size=9, color=FAINT)


def card(slide, x, y, w, h, head, body, head_color=CYAN, body_size=12):
    rect(slide, x, y, w, h, PANEL, rounded=True)
    rect(slide, x, y, Inches(0.06), h, head_color, rounded=False)
    pad = Inches(0.22)
    text(slide, x + pad, y + Inches(0.16), w - pad * 2, Inches(0.35),
         head, size=14, color=head_color, bold=True, font=TITLE_FONT)
    text(slide, x + pad, y + Inches(0.58), w - pad * 2, h - Inches(0.74),
         body, size=body_size, color=MUTED, line_spacing=1.12)


def chip(slide, x, y, w, label, color=LIME):
    c = rect(slide, x, y, w, Inches(0.34), PANEL_2, rounded=True, line_color=color, line_w=Pt(1))
    try:
        c.adjustments[0] = 0.5
    except Exception:
        pass
    text(slide, x, y + Inches(0.045), w, Inches(0.26), label, size=11, color=color,
         bold=True, align=PP_ALIGN.CENTER)


def arrow(slide, x, y, w, h=Inches(0.26)):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.fill.solid()
    a.fill.fore_color.rgb = CYAN
    _no_line(a)
    return a


def style_table(tbl, header_fill=PANEL_2, body_fill=PANEL, size=11, header_size=11):
    tbl_pr = tbl._tbl
    # remove default banding style
    for el in tbl_pr.findall(qn("a:tblPr")):
        el.set("bandRow", "0")
        el.set("firstRow", "1")
    for ri, row in enumerate(tbl.rows):
        for cell in row.cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if ri == 0 else body_fill
            cell.margin_left = cell.margin_right = Inches(0.1)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(header_size if ri == 0 else size)
                    run.font.name = BODY_FONT
                    run.font.bold = ri == 0
                    run.font.color.rgb = CYAN if ri == 0 else MUTED


def make_table(slide, x, y, w, h, data, col_widths=None, size=11):
    rows, cols = len(data), len(data[0])
    gfx = slide.shapes.add_table(rows, cols, x, y, w, h)
    tbl = gfx.table
    if col_widths:
        total = sum(col_widths)
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = Emu(int(w * cw / total))
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            tbl.cell(ri, ci).text = str(val)
    style_table(tbl, size=size)
    return tbl


# ================================================================ SLIDE 1 — Team Details
s = add_slide()
# hero accents
rect(s, 0, 0, SW, Inches(0.09), CYAN)
rect(s, Inches(10.4), Inches(0.9), Inches(2.4), Inches(0.055), LIME)

text(s, Inches(0.75), Inches(1.05), Inches(9), Inches(0.35),
     "SEMICON INDIA HACKATHON 2026  ·  IDEA SUBMISSION", size=13, color=CYAN, bold=True)
text(s, Inches(0.75), Inches(1.45), Inches(11.8), Inches(1.6),
     [("SpectraRestore", {"size": 54, "bold": True, "color": WHITE, "font": TITLE_FONT}),
      ("Fast Restoration of Semiconductor Inspection Images", {"size": 22, "color": MUTED})])
chip(s, Inches(0.75), Inches(3.15), Inches(3.3), "KLA · TRACK 1 · PS01", color=LIME)
chip(s, Inches(4.25), Inches(3.15), Inches(4.4), "AI-BASED IMAGE RESTORATION", color=CYAN)

text(s, Inches(0.75), Inches(3.85), Inches(6), Inches(0.3),
     "TEAM <TEAM NAME>  ·  <COLLEGE NAME>", size=16, color=WHITE, bold=True, font=TITLE_FONT)

make_table(
    s, Inches(0.75), Inches(4.35), Inches(11.8), Inches(2.1),
    [["Member", "Role", "Branch / Year", "Contact"],
     ["<Member 1>", "Team Lead · Model architecture & training", "<Branch, Year>", "<email / phone>"],
     ["<Member 2>", "Data pipeline & augmentation", "<Branch, Year>", "<email>"],
     ["<Member 3>", "Evaluation, metrics & benchmarking", "<Branch, Year>", "<email>"],
     ["<Member 4>", "Repo, documentation & presentation", "<Branch, Year>", "<email>"]],
    col_widths=[2.4, 4.6, 2.4, 2.4], size=12)

text(s, Inches(0.75), Inches(6.85), Inches(11.8), Inches(0.3),
     "Submission: TeamName_KLA_PS01.pdf   ·   GitHub repository linked on Slide 8",
     size=10, color=FAINT)

# ================================================================ SLIDE 2 — Problem Statement Addressed
s = add_slide()
header(s, "Slide 2 · Problem Statement Addressed", "AI-Based Restoration of Degraded Images", 2)

text(s, Inches(0.55), Inches(1.62), Inches(6.1), Inches(2.5),
     [("Why this matters", {"size": 15, "bold": True, "color": LIME, "font": TITLE_FONT}),
      ("Inspection tools capture microscopic grayscale images of chip structures to catch "
       "defects. A single pixel of noise or one lost detail can hide a defect — and a missed "
       "defect is a failed chip.", {"size": 13, "color": MUTED}),
      ("Today, engineers simply tolerate degraded images. Learned restoration upgrades every "
       "inspection tool's image quality in software — better defect detectability with zero "
       "new hardware, directly serving yield.", {"size": 13, "color": MUTED})],
     line_spacing=1.15, space_after=8)

text(s, Inches(0.55), Inches(4.15), Inches(6.1), Inches(2.6),
     [("The task", {"size": 15, "bold": True, "color": LIME, "font": TITLE_FONT}),
      ("Learn one model:  degraded 128² / 256² input  →  clean 2× output in [0, 1].",
       {"size": 13, "color": WHITE, "bold": True}),
      ("Judged on SSIM · pSNR · LPIPS, out-of-distribution robustness, and inference "
       "time on an H100 GPU.", {"size": 13, "color": MUTED}),
      ("KLA's own guidance: degradations \u201cmay appear in any order\u201d — so we invert the "
       "mapping from pairs instead of modelling the degradation chain.", {"size": 12, "color": FAINT, "italic": True})],
     line_spacing=1.15, space_after=8)

card(s, Inches(7.0), Inches(1.62), Inches(5.75), Inches(1.52), "1 · Speckle noise",
     "Multiplicative grain that pushes pixel values OUTSIDE [0, 1] while ground truth stays "
     "in [0, 1] — \u201ca feature, not a bug\u201d per KLA. Must be removed without blurring real structure.")
card(s, Inches(7.0), Inches(3.30), Inches(5.75), Inches(1.52), "2 · Additive Gaussian noise",
     "Softens edges and washes out fine structure. Sharpness must return without ringing "
     "or artificial patterns.")
card(s, Inches(7.0), Inches(4.98), Inches(5.75), Inches(1.52), "3 · 2× downsampling",
     "512² arrives as 256² (or 256² as 128²). Fine periodic detail is destroyed and must be "
     "reconstructed at 2× resolution.", head_color=LIME)

# ================================================================ SLIDE 3 — Idea Description
s = add_slide()
header(s, "Slide 3 · Idea Description", "One fast network, joint denoise + 2× super-resolution", 3)

rect(s, Inches(0.55), Inches(1.62), Inches(12.2), Inches(0.95), PANEL_2, rounded=True)
text(s, Inches(0.85), Inches(1.78), Inches(11.6), Inches(0.65),
     "A single NAFNet-based image-to-image network restores all three degradations in one forward pass — "
     "trained with a metric-matched loss (Charbonnier + SSIM + FFT + LPIPS) and degradation augmentation "
     "for out-of-distribution robustness.",
     size=14, color=WHITE, bold=True, line_spacing=1.15)

card(s, Inches(0.55), Inches(2.85), Inches(3.93), Inches(2.05), "Joint model, not a pipeline",
     "Denoise-then-upscale propagates errors between stages and doubles latency. One network "
     "shares features across denoising and detail reconstruction — faster AND better.")
card(s, Inches(4.68), Inches(2.85), Inches(3.93), Inches(2.05), "Regression, not generation",
     "No GAN, no diffusion. Generative models hallucinate texture — a hallucinated edge can fake "
     "or hide a defect. Regression directly optimises the announced metrics in a single pass.")
card(s, Inches(8.81), Inches(2.85), Inches(3.93), Inches(2.05), "Compute where it's cheap",
     "All heavy processing runs at LOW input resolution (~4× cheaper than at 2×); upsampling "
     "happens only at the very end. This is why the design wins the H100 timing benchmark.",
     head_color=LIME)

make_table(
    s, Inches(0.55), Inches(5.15), Inches(12.2), Inches(1.65),
    [["Degradation", "How the design handles it"],
     ["Speckle (out-of-range values)", "No input clipping; per-image standardisation absorbs the extended range; residual learning preserves true structure"],
     ["Gaussian noise", "Denoising backbone + SSIM & FFT loss terms restore edge sharpness without ringing"],
     ["2× resolution loss", "PixelShuffle sub-pixel tail reconstructs detail; FFT loss supervises the periodic spectrum"]],
    col_widths=[3.2, 9.0], size=11)

# ================================================================ SLIDE 4 — Proposed Solution
s = add_slide()
header(s, "Slide 4 · Proposed Solution", "Architecture — NAFNet-SR2× and the training recipe", 4)

# pipeline diagram
py = Inches(1.75)
ph = Inches(1.15)
boxes = [
    ("Degraded input", "1 × H × W\nvalues may exceed [0,1]", PANEL_2, MUTED),
    ("Standardise", "per-image\n(x − μ) / σ", PANEL, CYAN),
    ("NAFNet U-Net", "enc [2,2,4,8] · mid 12 · dec [2,2,2,2]\nwidth 32 · at input resolution", PANEL, CYAN),
    ("PixelShuffle ×2", "sub-pixel conv\nlearned upsampling", PANEL, CYAN),
    ("Restored output", "1 × 2H × 2W\nclamped to [0,1] at save", PANEL_2, LIME),
]
widths = [1.95, 1.75, 3.30, 1.95, 2.05]
gap = 0.28
x = 0.55
for (head_t, body_t, fill, hcol), bw in zip(boxes, widths):
    rect(s, Inches(x), py, Inches(bw), ph, fill, rounded=True)
    text(s, Inches(x + 0.12), py + Inches(0.10), Inches(bw - 0.24), Inches(0.3),
         head_t, size=13, bold=True, color=hcol, font=TITLE_FONT, align=PP_ALIGN.CENTER)
    text(s, Inches(x + 0.10), py + Inches(0.45), Inches(bw - 0.20), ph - Inches(0.5),
         body_t.split("\n"), size=10, color=MUTED, align=PP_ALIGN.CENTER, space_after=1)
    if bw != widths[-1] or x + bw < 12:
        pass
    x += bw + gap
# arrows between boxes
ax = 0.55
for bw in widths[:-1]:
    arrow(s, Inches(ax + bw + 0.015), py + Inches(0.45), Inches(gap - 0.03))
    ax += bw + gap

# global residual skip annotation
text(s, Inches(0.55), py + Inches(1.28), Inches(12.2), Inches(0.35),
     "+  global residual skip: noisy bilinear-upsampled input (absolute intensity) + learned correction. "
     "Input-only standardisation — output stays in GT space (no degraded mean/std re-applied).",
     size=11, color=LIME, align=PP_ALIGN.CENTER)

# training recipe cards
card(s, Inches(0.55), Inches(3.75), Inches(6.0), Inches(1.5), "Loss — matched to the metrics",
     "1.00 Charbonnier (pSNR proxy)  +  0.20 (1−SSIM)  +  0.05 FFT-L1 (periodic detail)  +  "
     "0.10 LPIPS enabled after 20% warm-up (prevents early artifacts). No adversarial loss.")
card(s, Inches(6.75), Inches(3.75), Inches(6.0), Inches(1.5), "Optimisation",
     "AdamW · LR 3e-4 → cosine to 1e-6 · ~200k iters · bf16 mixed precision · EMA 0.999 "
     "(EMA checkpoint ships) · checkpoint chosen by combined val SSIM + pSNR + (1−LPIPS).")
card(s, Inches(0.55), Inches(5.40), Inches(6.0), Inches(1.5), "Data & augmentation",
     "Held-out val/ split · aligned random crops (128² in / 256² GT) · flips + 90° "
     "rotations only · degradation aug: extra speckle / Gaussian / mild blur (p≈0.3).",
     head_color=LIME)
card(s, Inches(6.75), Inches(5.40), Inches(6.0), Inches(1.5), "Inference (what KLA benchmarks)",
     "FP16 + channels_last, single forward pass · ~29.2M params (width 32 default) · "
     "measure latency on target GPU · fast (~15M) preset if timing dominates · TTA off by default.",
     head_color=LIME)

# ================================================================ SLIDE 5 — Innovation & Uniqueness
s = add_slide()
header(s, "Slide 5 · Innovation & Uniqueness", "What makes this approach different", 5)

card(s, Inches(0.55), Inches(1.70), Inches(6.0), Inches(2.30), "1 · Degradation augmentation",
     "With p≈0.3 we inject EXTRA randomised speckle, Gaussian, or mild blur into already-degraded "
     "inputs (targets unchanged). Aligns with KLA webinar factors and widens the degradation family. "
     "Ablate with vs without augmentation. True OOD may also be new structures.",
     head_color=LIME)
card(s, Inches(6.75), Inches(1.70), Inches(6.0), Inches(2.30), "2 · Metric-aligned composite loss",
     "Charbonnier (pSNR proxy) + SSIM + FFT-L1 + LPIPS (after warm-up). These are the metrics we "
     "report on slides; H100 quality scoring is not fully specified. FFT term targets periodic "
     "semiconductor structure that plain L1 tends to blur.")
card(s, Inches(0.55), Inches(4.20), Inches(6.0), Inches(2.30), "3 · Input-only standardisation",
     "Per-image standardisation on the INPUT only absorbs out-of-range speckle and source shifts. "
     "We do NOT re-apply degraded mean/std to the output — predictions stay in absolute GT space "
     "([0,1]), as required by KLA range notes.")
card(s, Inches(6.75), Inches(4.20), Inches(6.0), Inches(2.30), "4 · Speed as a design constraint",
     "Attention-free NAFNet backbone, compute at low resolution, FP16 inference. Size presets: "
     "default ~29.2M, fast ~15M. Report a measured quality-vs-latency curve. Regression-only so "
     "the model cannot hallucinate defects — an argument for metrology judges.", head_color=LIME)

# ================================================================ SLIDE 6 — Results
s = add_slide()
header(s, "Slide 6 · Results", "Restoration quality on the held-out validation split", 6)

make_table(
    s, Inches(0.55), Inches(1.70), Inches(7.3), Inches(1.85),
    [["Metric", "Degraded input (baseline)", "Ours (restored)"],
     ["SSIM  ↑", "<baseline>", "<result>"],
     ["pSNR (dB)  ↑", "<baseline>", "<result>"],
     ["LPIPS  ↓", "<baseline>", "<result>"],
     ["Inference (ms/img, FP16)", "—", "<result>"]],
    col_widths=[3.0, 2.6, 2.2], size=12)

card(s, Inches(8.05), Inches(1.70), Inches(4.7), Inches(1.85), "OOD ablation",
     "Same table computed with vs without degradation augmentation on shifted-noise data — "
     "the number that proves generalisation.\n<fill after training>", body_size=11)

# before/after visual placeholders
text(s, Inches(0.55), Inches(3.80), Inches(12.2), Inches(0.3),
     "Visual evidence — degraded input → our restoration → ground truth (with |error| heat-map)",
     size=13, color=LIME, bold=True, font=TITLE_FONT)
labels = ["Degraded input", "Our restoration", "Ground truth", "|error| map"]
for i, lab in enumerate(labels):
    bx = Inches(0.55 + i * 3.13)
    ph_rect = rect(s, bx, Inches(4.20), Inches(2.93), Inches(2.15), PANEL_2, rounded=True,
                   line_color=FAINT, line_w=Pt(0.75))
    text(s, bx, Inches(5.05), Inches(2.93), Inches(0.4), "<image>", size=12, color=FAINT,
         align=PP_ALIGN.CENTER)
    text(s, bx, Inches(6.42), Inches(2.93), Inches(0.3), lab, size=11, color=MUTED,
         align=PP_ALIGN.CENTER, bold=True)

text(s, Inches(0.55), Inches(6.80), Inches(12.2), Inches(0.25),
     "Include at least one hard case (fine periodic array) and 2–3 zoom levels — confusion-free evidence as the brief demands.",
     size=10, color=FAINT)

# ================================================================ SLIDE 7 — Technology & Feasibility
s = add_slide()
header(s, "Slide 7 · Technology & Feasibility", "Standard, reproducible, student-compute friendly", 7)

make_table(
    s, Inches(0.55), Inches(1.70), Inches(6.4), Inches(4.4),
    [["Item", "Choice"],
     ["Framework", "PyTorch 2.x + lpips"],
     ["Model size", "~29.2M default (w32) · ~15M fast · ~65M large"],
     ["Training hardware", "Single A100 / L4 / T4-class cloud GPU"],
     ["Training time", "~200k iters (batch sized to GPU; Colab T4: batch 4–8)"],
     ["Inference", "FP16, single pass · report measured ms/image"],
     ["Data handling", "Native range preserved; png / tif / npy"],
     ["Reproducibility", "Fixed seeds · logged configs · pip freeze"]],
    col_widths=[2.1, 4.3], size=11)

card(s, Inches(7.25), Inches(1.70), Inches(5.5), Inches(2.0), "Why it's feasible",
     "Everything is standard and free of exotic dependencies. The model is deliberately small "
     "enough to train on student compute (Colab Pro / Kaggle), and a working checkpoint appears "
     "within hours — leaving time for the OOD ablation and visual assets.")
card(s, Inches(7.25), Inches(3.90), Inches(5.5), Inches(2.2), "Risk → mitigation",
     "Timing weighted heavily → switch to fast (~15M) preset\n"
     "OOD / blur variants → degradation aug (noise + mild blur)\n"
     "Format surprises → loader auto-detects; verify on real data first\n"
     "Filename scoring → evaluate.py writes same names as inputs",
     head_color=LIME, body_size=11)

# ================================================================ SLIDE 8 — GitHub & Video Link
s = add_slide()
header(s, "Slide 8 · GitHub & Video Link", "Repository & demo", 8)

rect(s, Inches(0.55), Inches(1.70), Inches(12.2), Inches(1.0), PANEL_2, rounded=True)
text(s, Inches(0.85), Inches(1.86), Inches(11.6), Inches(0.7),
     [("GitHub (mandatory):   github.com/<org>/<repo>", {"size": 16, "bold": True, "color": CYAN}),
      ("Demo video (optional):   <link — screen capture of evaluate.py restoring the test set>",
       {"size": 13, "color": MUTED})], space_after=4)

text(s, Inches(0.55), Inches(3.00), Inches(12.2), Inches(0.3),
     "Repository contents — exactly per the KLA checklist", size=14, color=LIME, bold=True,
     font=TITLE_FONT)

make_table(
    s, Inches(0.55), Inches(3.40), Inches(12.2), Inches(2.9),
    [["#", "Item", "What we ship"],
     ["1", "README.md", "Clone → install → download weights → run inference in ≤5 commands, no contact needed"],
     ["2", "evaluate.py (standalone)", "Takes --input_dir and --output_dir · same output filenames as inputs · auto-loads weights · zero manual edits · fresh-machine tested"],
     ["3", "train.py", "Reproduces the full training run from scratch (config + seeds)"],
     ["4", "Trained weights", "EMA checkpoint (.pt) via Git LFS / Drive / HuggingFace"],
     ["5", "outputs/", "Actual restored test images produced by our model"],
     ["6", "requirements.txt", "Full pip freeze from the training environment"]],
    col_widths=[0.5, 2.6, 9.1], size=11)

text(s, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.4),
     "The evaluation script is the critical file — KLA runs it AS-IS on the H100 to score quality and speed. "
     "It is tested on a clean machine before submission.",
     size=11, color=DANGER, bold=True)

# ================================================================ SLIDE 9 — References
s = add_slide()
header(s, "Slide 9 · References", "References", 9)

refs = [
    "1.  Chen, L. et al. — Simple Baselines for Image Restoration (NAFNet), ECCV 2022.",
    "2.  Shi, W. et al. — Real-Time Single Image and Video Super-Resolution Using an Efficient Sub-Pixel Convolutional Neural Network (PixelShuffle), CVPR 2016.",
    "3.  Zhang, R. et al. — The Unreasonable Effectiveness of Deep Features as a Perceptual Metric (LPIPS), CVPR 2018.",
    "4.  Zhai, L. et al. — A Comprehensive Review of Deep Learning-Based Real-World Image Restoration, IEEE Access 11, 2023.",
    "5.  Kumar, T. et al. — Image Data Augmentation Approaches: A Comprehensive Survey and Future Directions, IEEE Access 12, 2024.",
    "6.  Terven, J. et al. — A Comprehensive Survey of Loss Functions and Metrics in Deep Learning, Artificial Intelligence Review 58, 195, 2025.",
    "7.  Monga, V. et al. — Algorithm Unrolling: Interpretable, Efficient Deep Learning for Signal and Image Processing, IEEE Signal Processing Magazine 38(2), 2021.",
    "8.  Wang, Z. et al. — Image Quality Assessment: From Error Visibility to Structural Similarity (SSIM), IEEE Transactions on Image Processing, 2004.",
]
text(s, Inches(0.55), Inches(1.80), Inches(12.2), Inches(4.4),
     refs, size=13, color=MUTED, line_spacing=1.25, space_after=10)

text(s, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.5),
     "References 4–7 are the papers cited in KLA's own problem-statement webinar deck — our augmentation and loss "
     "choices are justified against them.",
     size=11, color=LIME)

# ---------------------------------------------------------------- save
out = "/Users/charnsuresh/Documents/Semicon/slides/TeamName_KLA_PS01.pptx"
prs.save(out)
print("saved", out)
